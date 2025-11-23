import os
import tempfile

import streamlit as st
from pptx import Presentation
from openai import OpenAI

# -------------------------
# CONFIG STREAMLIT
# -------------------------
st.set_page_config(
    page_title="Transcription & Résumé audio médical",
    page_icon="🎧",
    layout="wide"
)

# Bandeau haut personnalisé
st.markdown(
    "<h3 style='text-align: center; color: #005b96;'>créé avec amour par ton fils chéri &lt;3</h3>",
    unsafe_allow_html=True
)

st.title("🎧 Lilo & Mamati – Application de transcription audio")
st.write(
    "Lilo, pourrais-tu déposer ici un audio de conférence médicale (staff, cours, transmission…) "
    "puis choisir ce que tu souhaites : **Transcription complète**, **Résumé & points clés**, "
    "ou **Slides PowerPoint**. Pensée spéciale pour Mamati 💙."
)

# -------------------------
# CONFIG OPENAI (API KEY VIA SECRETS)
# -------------------------
if "OPENAI_API_KEY" not in st.secrets:
    st.error("⚠️ Clé API OpenAI manquante. Ajoute-la dans les secrets Streamlit (OPENAI_API_KEY).")
    st.stop()

client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

# -------------------------
# FONCTIONS MÉTIERS
# -------------------------
def transcribe_audio(uploaded_file) -> str:
    """
    Transcrit un fichier audio en texte avec le modèle 'whisper-1'
    (permet des durées plus longues que gpt-4o-mini-transcribe).
    """
    suffix = os.path.splitext(uploaded_file.name)[1] or ".mp3"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.read())
        tmp_path = tmp.name

    try:
        with open(tmp_path, "rb") as audio_file:
            transcription = client.audio.transcriptions.create(
                model="whisper-1",   # ✅ modèle sans limite stricte à 1400s
                file=audio_file,
                language="fr",       # français
            )
        text = transcription.text
    finally:
        # On nettoie le fichier temporaire
        if os.path.exists(tmp_path):
            os.remove(tmp_path)

    return text


def summarize_text(transcript: str) -> str:
    """Produit un résumé structuré de la transcription."""
    prompt = f"""
Tu es un médecin spécialiste qui résume des conférences médicales pour Lilo et Mamati.

Écris un résumé clair et structuré de la conférence ci-dessous.

Contraintes :
- En français.
- Commence par un résumé global en 5–10 lignes.
- Puis une section "Points clés" sous forme de bullet points.
- Puis une section "Implications pratiques pour la clinique" si pertinent (bullet points).
- Style : concis, pédagogique, sans phrases inutiles.

Transcription :
\"\"\"{transcript}\"\"\"
"""

    response = client.responses.create(
        model="gpt-5-nano",
        input=[
            {
                "role": "user",
                "content": prompt
            }
        ],
    )

    return response.output_text


def generate_slides_markdown(transcript: str) -> str:
    """
    Demande au modèle une structure de diaporama en Markdown.
    """
    prompt = f"""
À partir de cette transcription d'une conférence médicale, propose une structure de diaporama (PowerPoint) pour Lilo et Mamati, en français.

Contraintes :
- Entre 5 et 10 diapositives.
- Format STRICT en Markdown comme ci-dessous :
  # Titre de la présentation
  ## Slide 1 : Titre de la slide
  - Point 1
  - Point 2

  ## Slide 2 : Titre de la slide
  - Point 1
  - Point 2
  etc.

- La première diapositive doit être un titre général (sans puces).
- Les autres : objectifs, notions clés, physiopathologie, aspects cliniques, traitement, messages à retenir, conclusion.

Transcription :
\"\"\"{transcript}\"\"\"
"""

    response = client.responses.create(
        model="gpt-5-nano",
        input=[
            {
                "role": "user",
                "content": prompt
            }
        ],
    )

    return response.output_text


def markdown_to_pptx(md: str, output_path: str):
    """
    Transforme une structure de slides en Markdown en un fichier PPTX simple.
    """
    prs = Presentation()
    lines = [l.strip() for l in md.splitlines() if l.strip()]

    bullet_frame = None

    for line in lines:
        # Titre principal "# ..."
        if line.startswith("# ") and not line.startswith("##"):
            title_text = line[2:].strip()
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Titre
            slide.shapes.title.text = title_text
            continue

        # Nouvelle slide "## Slide X : Titre"
        if line.startswith("## "):
            slide_title = line[3:].strip()
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Titre + contenu
            slide.shapes.title.text = slide_title
            body = slide.placeholders[1]
            bullet_frame = body.text_frame
            bullet_frame.clear()
            continue

        # Bullet "- point"
        if line.startswith("- "):
            bullet_text = line[2:].strip()
            if bullet_frame is not None:
                if not bullet_frame.text:
                    bullet_frame.text = bullet_text
                else:
                    p = bullet_frame.add_paragraph()
                    p.text = bullet_text
            continue

    prs.save(output_path)


# -------------------------
# UI STREAMLIT
# -------------------------
uploaded_file = st.file_uploader(
    "Lilo, pourrais-tu déposer ici ton fichier audio (mp3, wav, m4a, mp4…) ?",
    type=["mp3", "wav", "m4a", "mp4"]
)

mode = st.radio(
    "Que veux-tu que l'application fasse pour toi, Lilo ?",
    [
        "Retranscription complète",
        "Résumé + points clés",
        "Génération de slides (PPTX)"
    ]
)

with st.expander("ℹ️ Conseils pour les audios longs (≈ 25–30 minutes)"):
    st.write(
        "- Lilo, pour des audios longs, privilégie si possible un format compressé (mp3).\n"
        "- Le traitement se fait côté OpenAI, donc même si l'audio est un peu long, "
        "l'application restera fluide pour toi et Mamati.\n"
        "- Si un jour un fichier est vraiment très long, on pourra envisager un découpage automatique."
    )

if uploaded_file is not None:
    st.audio(uploaded_file, format="audio/mp3")
    st.success("Merci Lilo 💙, l'audio est bien déposé. Choisis ce que tu veux en faire, puis lance le traitement.")

    if st.button("🚀 Lancer le traitement", type="primary"):
        try:
            with st.spinner("Lilo, je transcris l'audio pour toi…"):
                transcript = transcribe_audio(uploaded_file)

            if mode == "Retranscription complète":
                st.subheader("📝 Transcription complète")
                st.text_area("Texte transcrit", transcript, height=400)

                txt_path = "transcription.txt"
                with open(txt_path, "w", encoding="utf-8") as f:
                    f.write(transcript)

                with open(txt_path, "rb") as f:
                    st.download_button(
                        "📥 Télécharger la transcription (.txt)",
                        data=f,
                        file_name="transcription.txt",
                        mime="text/plain"
                    )

            elif mode == "Résumé + points clés":
                with st.spinner("Lilo, je prépare le résumé et les points clés…"):
                    summary = summarize_text(transcript)

                st.subheader("🧾 Résumé & points clés pour Lilo et Mamati")
                st.markdown(summary)

                txt_path = "resume_points_cles.txt"
                with open(txt_path, "w", encoding="utf-8") as f:
                    f.write(summary)

                with open(txt_path, "rb") as f:
                    st.download_button(
                        "📥 Télécharger le résumé (.txt)",
                        data=f,
                        file_name="resume_points_cles.txt",
                        mime="text/plain"
                    )

            elif mode == "Génération de slides (PPTX)":
                with st.spinner("Lilo, je génère la structure des slides pour ta conférence…"):
                    slides_md = generate_slides_markdown(transcript)

                st.subheader("📑 Structure des slides (Markdown généré)")
                st.markdown(slides_md)

                pptx_path = "conference_medicale_lilo_mamati.pptx"
                with st.spinner("Je crée le fichier PowerPoint pour toi, Lilo…"):
                    markdown_to_pptx(slides_md, pptx_path)

                with open(pptx_path, "rb") as f:
                    st.download_button(
                        "📥 Télécharger le diaporama (.pptx)",
                        data=f,
                        file_name="conference_medicale_lilo_mamati.pptx",
                        mime=(
                            "application/"
                            "vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    )

        except Exception as e:
            st.error(f"❌ Erreur lors du traitement : {e}")
            st.info("Si l'erreur persiste, envoie-moi un screen et on adaptera ensemble 💙.")
else:
    st.info("Lilo, pourrais-tu déposer un fichier audio pour commencer ?")